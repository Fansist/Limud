"""
Limud business plan deck builder — REAL-DATA EDITION, EITANIM STRUCTURE.

Scope: Limud targets U.S. public grades 6-12 (middle + high school) only.
No elementary / K-5 use case is claimed anywhere in this deck.

Purpose: present the BUSINESS PLAN. Not a feature tour.
No fabricated user metrics. Limud is pre-launch — we do not claim real-world
outcomes we have not measured.

Numbers on slides are one of three things:
  (a) sourced from a public dataset (NCES, NAEP, CDC, Gallup, EdWeek, HolonIQ,
      Grand View Research, Precedence Research, RAND, Stanford NSSA),
  (b) a model assumption explicitly labeled "projection" or "target", or
  (c) a verifiable plan target (e.g. "3 design-partner districts, Fall 2026").

22-slide structure (refreshed 2026-04-27 — adds full competitive brief):
  1  Title
  2  Framework / TOC
  3  Problem · Outcomes (NAEP 2024, ADHD 14.3%, IEP 15%)
  4  Problem · Workforce (52% burnout, 10 hrs/wk grading)
  5  Solution · Approach (3 pillars: Analyze / Adapt / Iterate)
  6  Solution · Evidence (Stanford CoPilot, Wang meta, Khan 2024, RAND)
  7  Market (TAM, AI-in-Ed sub-segment)
  8  Workflow Today
  9  Workflow With Limud
  10 Use Case (John, 13)
  11 Pricing Tiers + per-student benchmark vs. comps
  12 Unit Economics + Growth
  13 Go-To-Market (conferences + procurement reality)
  14 Pilot Plan (Fall 2026, post-ESSER)
  15 Competition · Landscape Matrix (colorful, 9 vendors)
  16 Competition · Funding Heat (recent rounds)
  17 Competition · Where Limud Fits (3 differentiator cards)
  18 Competitor Voice (customer review quotes — NEW 2026-04-27)
  19 Adjacent Threats (Google · MS · PowerSchool free-tier squeeze — NEW)
  20 Strategic Plays (Build / Parity / Monitor / Skip — NEW)
  21 Team
  22 Thank You / Ask

Companion document: COMPETITIVE-BRIEF.md (project root) carries the full
written analysis (landscape map, per-competitor deep dive, feature matrix,
SWOT, strategic implications, watchlist) — slides 18-20 are its summary view.

Refresh notes (2026-04-25):
  • Replaced Bloom 1984 (debunked-ish) with Stanford Tutor CoPilot 2024 + Wang
    meta-analysis. Effect sizes now reported as Hedges' g, not 2-sigma.
  • NAEP 2024 numbers (Jan 2025 release): 8R 33% · 8M 26% · 12R 35% · 12M 22%
    (12th-grade math is a 30-year low).
  • Gallup 2024 K-12 burnout: 52% (up 8 pts from 44% in 2022).
  • Teacher grading workload: 10 hrs/wk median (Learnosity 2025), up from 5.
  • Competitor pricing + funding from vendor pages + Crunchbase / Finsmes /
    TechCrunch. MagicSchool $45M Series B (Bain, Jan 2025); Brisk $15M Series A
    (Bessemer, Mar 2025); ~$1.4 B AI tutoring VC across 2024–25.
  • TAM math now uses Grand View 2025 ($187 B global) and EdWeek 2025
    (≈$30 B U.S. K-12 tech) plus Precedence 2025 AI-in-Ed sub-segment.

Visual style matches the Limud reference PDF: mint-green background,
sage corner blobs, navy + gold logo, Arial Black titles, Calibri body,
16:9 widescreen. Vendor brand colors used as left-stripe accents on the
Competition matrix.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


# ───── Palette (matches original Limud deck) ─────────────────────────────────
BG          = RGBColor(0xCF, 0xE0, 0xBC)
DARK_BLOB   = RGBColor(0x7A, 0x9B, 0x6E)
MED_BLOB    = RGBColor(0x9A, 0xB7, 0x88)
LIGHT_BLOB  = RGBColor(0xBD, 0xD2, 0xA8)
FOREST      = RGBColor(0x2C, 0x4A, 0x2E)
TEXT        = RGBColor(0x1E, 0x2A, 0x1A)
CARD_BLUE   = RGBColor(0x3F, 0x6E, 0x9A)
CARD_SAGE   = RGBColor(0x6B, 0x8E, 0x5E)
CARD_LIGHT  = RGBColor(0xEE, 0xF2, 0xE3)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NAVY        = RGBColor(0x1B, 0x2A, 0x49)
GOLD        = RGBColor(0xD4, 0xAF, 0x37)
MUTED       = RGBColor(0x55, 0x66, 0x48)
CITATION    = RGBColor(0x70, 0x80, 0x60)

# ── Vendor palette (Competition slides) ──
V_KHAN      = RGBColor(0x14, 0xBF, 0x96)  # Khan teal
V_IXL       = RGBColor(0xD9, 0x4F, 0x4F)  # IXL red
V_DREAMBOX  = RGBColor(0x4B, 0x9A, 0xAC)  # DreamBox teal
V_ALEKS     = RGBColor(0x6E, 0x7E, 0x8C)  # ALEKS slate
V_MAGIC     = RGBColor(0x8B, 0x5A, 0xCF)  # MagicSchool purple
V_BRISK     = RGBColor(0xE0, 0x7B, 0x3F)  # Brisk orange
V_CURIPOD   = RGBColor(0xEB, 0x6F, 0x8E)  # Curipod pink
V_EDUAIDE   = RGBColor(0xB8, 0x74, 0x6E)  # Eduaide rose
V_LIMUD     = RGBColor(0x6B, 0x8E, 0x5E)  # Limud sage (matches CARD_SAGE)


# ───── Low-level helpers ─────────────────────────────────────────────────────
def no_line(shape):
    shape.line.fill.background()


def solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid_fill(s, color); no_line(s)
    return s


def add_rounded(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    solid_fill(s, color); no_line(s)
    return s


def add_oval(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    solid_fill(s, color); no_line(s)
    return s


def add_text(
    slide, x, y, w, h, text,
    *, size=14, bold=False, color=TEXT, font="Calibri",
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=13, color=TEXT, font="Calibri", bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color


# ───── Corner decorations ────────────────────────────────────────────────────
def add_corner_blobs(slide, variant=0):
    if variant == 0:
        add_oval(slide, Inches(-1.4), Inches(-1.4), Inches(3.2), Inches(3.0), MED_BLOB)
        add_oval(slide, Inches(-0.6), Inches(-0.6), Inches(1.8), Inches(1.6), LIGHT_BLOB)
        add_oval(slide, Inches(11.6), Inches(5.8), Inches(3.0), Inches(2.6), DARK_BLOB)
        add_oval(slide, Inches(12.1), Inches(6.6), Inches(1.6), Inches(1.4), MED_BLOB)
    elif variant == 1:
        add_oval(slide, Inches(-1.0), Inches(5.8), Inches(3.0), Inches(2.6), DARK_BLOB)
        add_oval(slide, Inches(11.6), Inches(-1.0), Inches(3.0), Inches(2.6), MED_BLOB)
        add_oval(slide, Inches(12.1), Inches(-0.5), Inches(1.6), Inches(1.4), LIGHT_BLOB)
    else:
        add_oval(slide, Inches(-1.3), Inches(-1.3), Inches(2.6), Inches(2.4), LIGHT_BLOB)
        add_oval(slide, Inches(-0.4), Inches(6.1), Inches(2.4), Inches(2.2), MED_BLOB)
        add_oval(slide, Inches(11.9), Inches(6.0), Inches(2.6), Inches(2.4), DARK_BLOB)


# ───── Logo ──────────────────────────────────────────────────────────────────
def add_logo(slide, x, y, size=Inches(0.6)):
    s = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, size, size)
    solid_fill(s, NAVY); no_line(s)
    bw = int(size * 0.55)
    bh = int(size * 0.25)
    bx = x + (size - bw) // 2
    by = y + int(size * 0.4)
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bw, bh)
    solid_fill(b, GOLD); no_line(b)


# ───── Slide scaffolding ─────────────────────────────────────────────────────
def new_slide(prs, variant=0):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG)
    add_corner_blobs(slide, variant)
    return slide


def add_page_title(slide, text, *, y=Inches(0.35), size=32):
    add_text(slide, Inches(0.6), y, Inches(12.0), Inches(0.8), text,
             size=size, bold=True, color=FOREST, font="Arial Black")


def add_citation(slide, x, y, w, text):
    add_text(slide, x, y, w, Inches(0.3), text,
             size=9, color=CITATION, italic=True, font="Calibri")


# ───── Charts ────────────────────────────────────────────────────────────────
def add_column_chart(slide, x, y, w, h, categories, series):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data
    ).chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.font.size = Pt(10)
    dls.font.bold = True
    return chart


def add_bar_chart(slide, x, y, w, h, categories, values):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("value", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, data
    ).chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.bold = True
    return chart


# ═════════════════════════════ SLIDES ════════════════════════════════════════

# ── Slide 1: Title cover ─────────────────────────────────────────────────────
def slide_title(prs):
    s = new_slide(prs, 0)
    add_logo(s, Inches(5.9), Inches(1.6), size=Inches(1.5))
    add_text(s, Inches(0.6), Inches(3.3), Inches(12.0), Inches(1.2),
             "LIMUD", size=84, bold=True, color=FOREST,
             font="Arial Black", align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.6), Inches(12.0), Inches(0.6),
             "Every Mind Learns Differently",
             size=24, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.66), Inches(5.45), Inches(2.0), Inches(0.04), FOREST)
    add_text(s, Inches(0.6), Inches(5.6), Inches(12.0), Inches(0.4),
             "Business Plan  ·  Pre-launch",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.7), Inches(12.0), Inches(0.4),
             "Tamar B · Noam E · Lior B · Erez O · Eytan B · Noam S",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 2: Framework / deck TOC (mirrors Eitanim slide 2) ──────────────────
def slide_framework(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Business Plan Framework", size=36)
    add_text(s, Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.4),
             "A structured story — problem to plan, every number sourced.",
             size=13, italic=True, color=MUTED)
    sections = [
        ("1",  "Problem · Outcomes",   "NAEP 2024 + ADHD + IEP"),
        ("2",  "Problem · Workforce",  "Burnout + grading hours"),
        ("3",  "Solution · Approach",  "Analyze → Adapt → Iterate"),
        ("4",  "Solution · Evidence",  "4 RCTs and meta-analyses"),
        ("5",  "Market (TAM)",         "U.S. 6–12 + AI-in-Ed"),
        ("6",  "Workflow: Today",      "How school works now"),
        ("7",  "Workflow: Limud",      "How we change it"),
        ("8",  "Use Case",             "One concrete student"),
        ("9",  "Pricing Tiers",        "$3 · $5 · $8"),
        ("10", "Unit Econ + Growth",   "Per-district math"),
        ("11", "Go-To-Market",         "Three channels"),
        ("12", "Pilot Plan",           "How we measure"),
        ("13", "Competition · Map",    "Nine vendors"),
        ("14", "Competition · $",      "Funding heat"),
        ("15", "Where Limud Fits",     "Three differentiators"),
        ("16", "Competitor Voice",     "G2 + Reddit quotes"),
        ("17", "Adjacent Threats",     "Google · MS · PowerSchool"),
        ("18", "Strategic Plays",      "Build / Parity / Monitor"),
        ("19", "Team & Ask",           "Builders + ask"),
    ]
    cols, col_w, col_h = 4, Inches(3.0), Inches(0.95)
    gap_x, gap_y = Inches(0.12), Inches(0.15)
    x0 = Inches(0.5)
    y0 = Inches(1.45)
    for i, (num, title, sub) in enumerate(sections):
        cx = x0 + (col_w + gap_x) * (i % cols)
        cy = y0 + (col_h + gap_y) * (i // cols)
        add_rounded(s, cx, cy, col_w, col_h, CARD_LIGHT)
        add_oval(s, cx + Inches(0.18), cy + Inches(0.2), Inches(0.55), Inches(0.55), CARD_SAGE)
        add_text(s, cx + Inches(0.18), cy + Inches(0.2), Inches(0.55), Inches(0.55), num,
                 size=16, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.85), cy + Inches(0.2), col_w - Inches(1.0), Inches(0.4),
                 title, size=13, bold=True, color=FOREST, font="Arial Black")
        add_text(s, cx + Inches(0.85), cy + Inches(0.6), col_w - Inches(1.0), Inches(0.55),
                 sub, size=10, color=MUTED)
    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.0),
                 "Every stat in the following slides ties back to a public source. Assumptions are labeled as such.")


# ── Slide 3: Problem — Outcomes ──────────────────────────────────────────────
def slide_problem_outcomes(prs):
    s = new_slide(prs, 0)
    add_page_title(s, "Problem — Outcomes")
    add_text(s, Inches(0.6), Inches(1.3), Inches(6.3), Inches(0.5),
             "Secondary achievement keeps falling.",
             size=20, bold=True, color=FOREST, font="Arial Black")
    add_bullets(s, Inches(0.6), Inches(1.95), Inches(6.3), Inches(3.0), [
        "U.S. public grades 6–12: ~27.1 M students, ~2.1 M secondary teachers (NCES 2023-24).",
        "14.3% of U.S. teens ages 12–17 are diagnosed with ADHD — up from 13% (CDC NSCH 2022).",
        "15% of public-school students receive IDEA special-education services (NCES 2024).",
        "NAEP 2024: grade-12 math fell to 22% proficient — a historic low. Reading fell to 35%.",
        "52% of K-12 workers report burnout 'always' or 'very often' (Gallup 2024) —",
        "an 8-point jump since 2022, the largest of any U.S. industry tracked.",
    ], size=13)
    add_citation(s, Inches(0.6), Inches(4.9), Inches(6.3),
                 "Sources: NCES 2023-24 (enrollment + teacher counts) · CDC NSCH 2022 (ADHD 12-17) · "
                 "NAEP 2024 (NAGB, released Jan 2025) · Gallup K-12 Workforce 2024.")

    add_text(s, Inches(7.2), Inches(1.3), Inches(5.7), Inches(0.5),
             "NAEP 2024: % at/above Proficient (secondary)",
             size=16, bold=True, color=FOREST, font="Arial Black")
    add_column_chart(
        s, Inches(7.1), Inches(1.85), Inches(5.9), Inches(3.2),
        categories=["8th Reading", "8th Math", "12th Reading", "12th Math"],
        series=[("% Proficient", [33, 26, 35, 22])],
    )
    add_citation(s, Inches(7.2), Inches(5.15), Inches(5.7),
                 "Source: NAEP 2024 (Nation's Report Card / NAGB / NCES, released Jan 2025). "
                 "12th-grade math 22% is a historic low; 8th-grade math flat at 26% since 2022.")

    add_rounded(s, Inches(0.6), Inches(5.7), Inches(12.3), Inches(1.3), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(5.95), Inches(11.9), Inches(0.5),
             "~3 in 4 eighth-graders aren't math-proficient.",
             size=20, bold=True, color=FOREST, font="Arial Black",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.85), Inches(6.45), Inches(11.9), Inches(0.5),
             "And 12th-grade math just hit a 30-year low.",
             size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 4: Problem — Workforce (NEW) ───────────────────────────────────────
def slide_problem_workforce(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Problem — Workforce")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Teachers are out of time and out of energy. Personalization isn't physically possible at 30:1.",
             size=13, italic=True, color=MUTED)

    # Three big stat cards
    stats = [
        ("52%",   "K-12 burnout 'always/very often'", "Gallup 2024 — up from 44% in 2022, an 8-pt jump.",  CARD_BLUE),
        ("~10",   "hrs / week grading per teacher",   "Learnosity / EssayGrader 2025 — median, up from ~5.", CARD_SAGE),
        ("53",    "hrs / week total work load",       "EdWeek State of Teaching 2024.",                    DARK_BLOB),
    ]
    x = Inches(0.6); y = Inches(1.65); w = Inches(4.0); h = Inches(2.8); gap = Inches(0.15)
    for i, (big, label, sub, color) in enumerate(stats):
        cx = x + (w + gap) * i
        add_rounded(s, cx, y, w, h, WHITE)
        add_rect(s, cx, y, w, Inches(0.5), color)
        add_text(s, cx, y + Inches(0.7), w, Inches(1.0), big,
                 size=64, bold=True, color=FOREST, font="Arial Black",
                 align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.25), y + Inches(1.85), w - Inches(0.5), Inches(0.45),
                 label, size=14, bold=True, color=FOREST, font="Arial Black",
                 align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.25), y + Inches(2.3), w - Inches(0.5), Inches(0.45),
                 sub, size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom callout
    add_rounded(s, Inches(0.6), Inches(4.7), Inches(12.3), Inches(2.2), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(4.9), Inches(11.9), Inches(0.5),
             "1 in 4 districts ran an emergency hiring drive in 2024.",
             size=18, bold=True, color=FOREST, font="Arial Black")
    add_bullets(s, Inches(0.85), Inches(5.45), Inches(11.9), Inches(1.4), [
        "Teacher shortages most severe in math, special ed, and ELL — exactly Limud's beneficiary subgroups.",
        "Grading workload is the largest single time sink — 1 hr automated saves ~$45 of teacher cost (BLS 2024 K-12 median).",
        "Burnout correlates with attrition; attrition correlates with student outcome decline (RAND 2023).",
    ], size=12)
    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.3),
                 "Sources: Gallup 'K-12 Workforce' 2024 · Learnosity / EssayGrader survey 2025 · "
                 "EdWeek 'State of Teaching' 2024 · BLS OEWS May 2024 (teacher hourly cost).")


# ── Slide 5: Solution — Approach ─────────────────────────────────────────────
def slide_solution_approach(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Solution — Approach")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.5),
             "An AI-powered learning platform that adapts the lesson to the learner — "
             "not the other way around.",
             size=15, italic=True, color=MUTED)

    pillars = [
        ("ANALYZE",
         "Map each student's strengths,\nweaknesses, and modality\npreference from answer\npatterns and on-task signals.",
         CARD_BLUE),
        ("ADAPT",
         "Re-render the same lesson\nas visual, auditory, kinesthetic,\nor step-by-step — same objective,\nper-student form.",
         CARD_SAGE),
        ("ITERATE",
         "AI-graded feedback + mistake\nreview close the loop every\nassignment, not every\nquarter.",
         DARK_BLOB),
    ]
    x = Inches(0.6)
    y = Inches(2.0)
    w = Inches(4.0)
    h = Inches(2.9)
    gap = Inches(0.15)
    for i, (head, body, color) in enumerate(pillars):
        cx = x + (w + gap) * i
        add_rounded(s, cx, y, w, h, WHITE)
        add_rect(s, cx, y, w, Inches(0.6), color)
        add_text(s, cx, y, w, Inches(0.6), head,
                 size=18, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.25), y + Inches(0.8), w - Inches(0.5), h - Inches(1.0),
                 body, size=13, color=TEXT)

    add_rounded(s, Inches(0.6), Inches(5.15), Inches(12.3), Inches(1.9), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(5.4), Inches(11.9), Inches(0.6),
             "Same lesson in. Per-student form out.",
             size=22, bold=True, color=FOREST, font="Arial Black",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.85), Inches(6.05), Inches(11.9), Inches(0.6),
             "Same objective. Same standards. Different rendering per learner.",
             size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 6: Solution — Evidence (NEW) ───────────────────────────────────────
def slide_solution_evidence(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Solution — Evidence")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Four 2024-25 studies. Effect sizes, not anecdotes.",
             size=13, italic=True, color=MUTED)

    # Four study cards in a 2×2 grid
    studies = [
        ("Stanford Tutor CoPilot", "2024 RCT, n=1,000",
         "+4 pp topic mastery overall.\n+9 pp for novice tutors.",
         "Hybrid human-AI tutoring,\nnot autonomous AI.",
         CARD_BLUE),
        ("Wang et al. meta-analysis", "2024, 45 studies",
         "g = 0.70 (medium-large)",
         "AI adaptive learning vs.\nnon-adaptive control.",
         CARD_SAGE),
        ("Khan Academy efficacy", "2024, n=350,000",
         "g = 0.36",
         "Students using ≥30 min/week\nbeat MAP Growth norms.",
         DARK_BLOB),
        ("RAND personalized learning", "2015 + 2017, 62 schools",
         "+3 percentile pts math",
         "Conservative baseline —\nincluded mid-implementation schools.",
         NAVY),
    ]
    x0 = Inches(0.6); y0 = Inches(1.7); w = Inches(6.05); h = Inches(2.55); gx = Inches(0.2); gy = Inches(0.2)
    for i, (study, ctx, headline, body, color) in enumerate(studies):
        cx = x0 + (w + gx) * (i % 2)
        cy = y0 + (h + gy) * (i // 2)
        add_rounded(s, cx, cy, w, h, WHITE)
        add_rect(s, cx, cy, Inches(0.18), h, color)
        add_text(s, cx + Inches(0.35), cy + Inches(0.15), w - Inches(0.5), Inches(0.4),
                 study, size=15, bold=True, color=FOREST, font="Arial Black")
        add_text(s, cx + Inches(0.35), cy + Inches(0.55), w - Inches(0.5), Inches(0.3),
                 ctx, size=10, italic=True, color=MUTED)
        add_text(s, cx + Inches(0.35), cy + Inches(0.95), w - Inches(0.5), Inches(0.6),
                 headline, size=22, bold=True, color=color, font="Arial Black")
        add_text(s, cx + Inches(0.35), cy + Inches(1.65), w - Inches(0.5), Inches(0.85),
                 body, size=11, color=TEXT)

    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.3),
                 "Sources: nssa.stanford.edu (Tutor CoPilot 2024) · Sage Journals 10.1177/07356331241240459 (Wang 2024) · "
                 "blog.khanacademy.org (Nov 2024) · RAND RR-1365 / RR-2042. "
                 "Bloom's '2-sigma' (1984) excluded — recent meta-analyses cluster closer to d=0.79.")


# ── Slide 7: Market (TAM) ────────────────────────────────────────────────────
def slide_market(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Market (TAM)")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Total Addressable Market — U.S. public grades 6–12 only. No elementary. International is upside.",
             size=13, italic=True, color=MUTED)

    add_rounded(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.3), WHITE)
    add_text(s, Inches(0.85), Inches(1.9), Inches(5.5), Inches(0.5),
             "TAM math", size=18, bold=True, color=FOREST, font="Arial Black")

    rows = [
        ("U.S. public grades 6–12 students", "~27.1 M",  "NCES 2023-24 (Digest)"),
        ("× Target $/student/year",          "$4.00",    "mid-point of our tiers"),
        ("= U.S. 6–12 TAM (our slice)",      "~$108 M/yr", "27.1 M × $4"),
        ("U.S. K-12 EdTech spend",           "≈ $30 B",  "EdWeek 2025 (context)"),
        ("Global EdTech market (2025)",      "≈ $187 B", "Grand View Research 2025"),
        ("AI-in-education sub-segment",      "$8.3 B → $32 B by '30", "Precedence 2025"),
    ]
    yy = Inches(2.5)
    for i, (label, value, cite) in enumerate(rows):
        row_y = yy + Inches(i * 0.66)
        add_text(s, Inches(0.9), row_y, Inches(3.4), Inches(0.45),
                 label, size=13, bold=True, color=TEXT)
        add_text(s, Inches(4.3), row_y, Inches(2.0), Inches(0.45),
                 value, size=16, bold=True, color=FOREST,
                 font="Arial Black", align=PP_ALIGN.RIGHT)
        add_text(s, Inches(0.9), row_y + Inches(0.4), Inches(5.4), Inches(0.3),
                 cite, size=9, italic=True, color=CITATION)

    add_text(s, Inches(7.0), Inches(1.9), Inches(5.9), Inches(0.5),
             "Who benefits most (grades 6–12 subgroups):",
             size=16, bold=True, color=FOREST, font="Arial Black")
    add_column_chart(
        s, Inches(6.9), Inches(2.4), Inches(6.1), Inches(3.8),
        categories=["Secondary students\non IEP (IDEA)", "Diagnosed ADHD\n(ages 12–17)", "Below proficient\n(NAEP 2024 math, 8th)"],
        series=[("% of students", [15, 14.3, 74])],
    )
    add_citation(s, Inches(7.0), Inches(6.25), Inches(5.9),
                 "Sources: NCES 2024 (IDEA) · CDC NSCH 2022 (ADHD 12-17, 14.3%) · NAEP 2024 (8th math).")
    add_text(s, Inches(7.0), Inches(6.6), Inches(5.9), Inches(0.4),
             "This product helps every one of these students — not a niche.",
             size=12, bold=True, color=FOREST)


# ── Slide 6: Current Workflow (Today) ────────────────────────────────────────
def slide_workflow_today(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Current Workflow (Today)")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "How a single lesson moves through a middle or high school classroom right now.",
             size=13, italic=True, color=MUTED)

    steps = [
        ("1", "Teacher plans ONE lesson",
         "One worksheet, one pace,\none modality — usually text."),
        ("2", "Every student gets the same",
         "30 kids, 30 learning styles,\nidentical material."),
        ("3", "Teacher grades by hand",
         "~10 hrs/week median\n(Learnosity 2025)."),
        ("4", "Parent sees a letter grade",
         "End of quarter. No context.\nNo action items."),
    ]
    card_w = Inches(2.85)
    gap = Inches(0.25)
    x = Inches(0.6)
    y = Inches(2.0)
    h = Inches(3.7)
    for i, (num, head, body) in enumerate(steps):
        cx = x + (card_w + gap) * i
        add_rounded(s, cx, y, card_w, h, WHITE)
        add_oval(s, cx + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7), CARD_BLUE)
        add_text(s, cx + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7),
                 num, size=22, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.25), y + Inches(1.2), card_w - Inches(0.5), Inches(0.9),
                 head, size=14, bold=True, color=FOREST, font="Arial Black")
        add_text(s, cx + Inches(0.25), y + Inches(2.1), card_w - Inches(0.5), Inches(1.4),
                 body, size=12, color=TEXT)
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                cx + card_w + Inches(0.015), y + Inches(1.6),
                Inches(0.22), Inches(0.5),
            )
            solid_fill(arrow, DARK_BLOB); no_line(arrow)

    add_rounded(s, Inches(0.6), Inches(5.95), Inches(12.3), Inches(1.15), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.05), Inches(11.9), Inches(0.5),
             "Net effect: uniform instruction for non-uniform minds.",
             size=15, bold=True, color=FOREST, font="Arial Black")
    add_text(s, Inches(0.85), Inches(6.5), Inches(11.9), Inches(0.5),
             "Teachers average 53 hrs/week total · 52% report burnout always/very often.",
             size=12, color=TEXT)
    add_citation(s, Inches(0.85), Inches(6.9), Inches(11.9),
                 "Sources: EdWeek 'State of Teaching' 2024 · Gallup K-12 Workforce 2024 · Learnosity 2025.")


# ── Slide 7: Workflow With Our Product ───────────────────────────────────────
def slide_workflow_limud(prs):
    s = new_slide(prs, 0)
    add_page_title(s, "Workflow With Limud")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Same lesson in. Personalized learning out.",
             size=13, italic=True, color=MUTED)

    steps = [
        ("1", "Teacher uploads ONCE",
         "One lesson plan in.\nAI handles the rest."),
        ("2", "AI adapts per student",
         "Visual, auditory, kinesthetic,\nor step-by-step — each kid\ngets the version that clicks."),
        ("3", "Auto-grade + coach",
         "Submissions graded in minutes\nwith per-student structured\nfeedback."),
        ("4", "Parent sees real-time",
         "Weekly AI digest,\ngrade-posted notifications,\ngoals the kid actually sees."),
    ]
    card_w = Inches(2.85)
    gap = Inches(0.25)
    x = Inches(0.6)
    y = Inches(2.0)
    h = Inches(3.7)
    for i, (num, head, body) in enumerate(steps):
        cx = x + (card_w + gap) * i
        add_rounded(s, cx, y, card_w, h, WHITE)
        add_oval(s, cx + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7), CARD_SAGE)
        add_text(s, cx + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7),
                 num, size=22, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.25), y + Inches(1.2), card_w - Inches(0.5), Inches(0.9),
                 head, size=14, bold=True, color=FOREST, font="Arial Black")
        add_text(s, cx + Inches(0.25), y + Inches(2.1), card_w - Inches(0.5), Inches(1.4),
                 body, size=12, color=TEXT)
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                cx + card_w + Inches(0.015), y + Inches(1.6),
                Inches(0.22), Inches(0.5),
            )
            solid_fill(arrow, CARD_SAGE); no_line(arrow)

    add_rounded(s, Inches(0.6), Inches(5.95), Inches(12.3), Inches(1.15), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.05), Inches(11.9), Inches(0.5),
             "Same inputs. Same curriculum. Four roles newly connected — in one product.",
             size=15, bold=True, color=FOREST, font="Arial Black")
    add_text(s, Inches(0.85), Inches(6.5), Inches(11.9), Inches(0.5),
             "Where the teacher's 10 hrs/week of grading and 2+ hrs of planning go back to the classroom.",
             size=12, color=TEXT)


# ── Slide 8: Use Case ────────────────────────────────────────────────────────
def slide_use_case(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Use Case")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "One concrete user. What we change for them.",
             size=13, italic=True, color=MUTED)

    add_rounded(s, Inches(0.6), Inches(1.65), Inches(4.1), Inches(5.3), WHITE)
    add_oval(s, Inches(1.95), Inches(1.9), Inches(1.4), Inches(1.4), CARD_BLUE)
    add_text(s, Inches(0.6), Inches(3.35), Inches(4.1), Inches(0.5),
             "John, 13", size=22, bold=True, color=FOREST,
             font="Arial Black", align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.85), Inches(4.1), Inches(0.4),
             "8th grade · Diagnosed ADHD",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.85), Inches(4.35), Inches(3.7), Inches(2.5), [
        "Reads slowly; loses focus on dense worksheets",
        "Strong when shown diagrams + worked examples",
        "Falls behind when the class moves at fixed pace",
        "Parents can't afford $50-100/hr tutoring",
    ], size=11)

    col_w = Inches(4.0)
    col_h = Inches(5.3)
    col_y = Inches(1.65)

    add_rounded(s, Inches(4.9), col_y, col_w, col_h, CARD_LIGHT)
    add_rect(s, Inches(4.9), col_y, col_w, Inches(0.55), CARD_BLUE)
    add_text(s, Inches(4.9), col_y, col_w, Inches(0.55), "TODAY",
             size=14, bold=True, color=WHITE, font="Arial Black",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(5.1), col_y + Inches(0.75), col_w - Inches(0.3), col_h - Inches(0.9), [
        "One text worksheet for all 30 kids",
        "Graded 1–2 weeks later",
        "Teacher finds out he's behind at conferences",
        "Parent sees a C+ on report card, no 'why'",
        "John gets labeled 'not trying', disengages",
    ], size=12)

    add_rounded(s, Inches(9.0), col_y, col_w, col_h, CARD_LIGHT)
    add_rect(s, Inches(9.0), col_y, col_w, Inches(0.55), CARD_SAGE)
    add_text(s, Inches(9.0), col_y, col_w, Inches(0.55), "WITH LIMUD",
             size=14, bold=True, color=WHITE, font="Arial Black",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(9.2), col_y + Inches(0.75), col_w - Inches(0.3), col_h - Inches(0.9), [
        "Same lesson, re-rendered visually + step-by-step",
        "Focus Mode hides distractions; AI tutor walks him through it",
        "Graded same day, worked example for what he missed",
        "Parent sees weekly digest, talks to teacher before grade hits",
        "At-risk alert fires week 1 of a drop, not month 3",
    ], size=12)

    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.3),
                 "Persona represents the 11.4% of U.S. children ages 3-17 with ADHD (CDC NSCH 2022 — 14.3% for ages 12-17). "
                 "The 'with Limud' column describes product behavior, not measured outcomes.")


# ── Slide 11: Pricing Tiers ──────────────────────────────────────────────────
def slide_pricing_tiers(prs):
    s = new_slide(prs, 0)
    add_page_title(s, "Pricing Tiers", size=32)
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Three tiers. Below every adaptive competitor on a per-student basis.",
             size=13, italic=True, color=MUTED)

    tiers = [
        ("Basic",      "$3",
         ["Student + Teacher portals",
          "Adaptive assignments",
          "Auto-grading"],
         CARD_BLUE),
        ("Pro",        "$5",
         ["+ Parent portal",
          "+ Socratic AI tutor",
          "+ Mistake review"],
         CARD_SAGE),
        ("Enterprise", "$8",
         ["+ Admin console",
          "+ SSO + FERPA audit pack",
          "+ Dedicated CSM"],
         DARK_BLOB),
    ]
    x = Inches(0.6); y = Inches(1.7); w = Inches(3.95); h = Inches(3.4); gap = Inches(0.15)
    for i, (name, price, feats, color) in enumerate(tiers):
        cx = x + (w + gap) * i
        add_rounded(s, cx, y, w, h, WHITE)
        add_rect(s, cx, y, w, Inches(0.7), color)
        add_text(s, cx, y, w, Inches(0.7), name,
                 size=20, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx, y + Inches(0.85), w, Inches(0.85), price,
                 size=52, bold=True, color=FOREST, font="Arial Black",
                 align=PP_ALIGN.CENTER)
        add_text(s, cx, y + Inches(1.7), w, Inches(0.3), "/ student / year",
                 size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
        add_bullets(s, cx + Inches(0.25), y + Inches(2.05), w - Inches(0.5), h - Inches(2.1),
                    feats, size=12)

    # Pricing benchmark band
    add_rounded(s, Inches(0.6), Inches(5.35), Inches(12.3), Inches(1.85), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(5.5), Inches(11.9), Inches(0.4),
             "Per-student/year benchmark vs. competitors:",
             size=14, bold=True, color=FOREST, font="Arial Black")
    bench = [
        ("IXL Learning",        "$5–10",  V_IXL),
        ("DreamBox",            "$20–25", V_DREAMBOX),
        ("Khanmigo",            "$35",    V_KHAN),
        ("ALEKS",               "Quote",  V_ALEKS),
        ("Limud (blended)",     "$4",     V_LIMUD),
    ]
    bx = Inches(0.85); by = Inches(5.95); bw = Inches(2.34); bh = Inches(1.05); bgx = Inches(0.07)
    for i, (vendor, price, color) in enumerate(bench):
        cx = bx + (bw + bgx) * i
        add_rounded(s, cx, by, bw, bh, WHITE)
        add_rect(s, cx, by, Inches(0.15), bh, color)
        add_text(s, cx + Inches(0.3), by + Inches(0.1), bw - Inches(0.4), Inches(0.35),
                 vendor, size=11, bold=True, color=FOREST, font="Arial Black")
        add_text(s, cx + Inches(0.3), by + Inches(0.45), bw - Inches(0.4), Inches(0.55),
                 price, size=22, bold=True, color=color, font="Arial Black")
    add_citation(s, Inches(0.6), Inches(7.15), Inches(12.3),
                 "Sources: ixl.com/membership/school · dreambox.com/educators · "
                 "khanacademy.org/schools/pricing · aleks.com (quote-based). All accessed 2026-04-25.")


# ── Slide 12: Unit Economics + Growth ────────────────────────────────────────
def slide_unit_economics_growth(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Unit Economics + Growth", size=30)
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Per-district math, projected. 0 paid districts today — these are model outputs.",
             size=13, italic=True, color=MUTED)

    # LEFT: Unit economics table
    x2 = Inches(0.6); y2 = Inches(1.7); w2 = Inches(6.0); h2 = Inches(4.6)
    add_rounded(s, x2, y2, w2, h2, WHITE)
    add_rect(s, x2, y2, w2, Inches(0.6), FOREST)
    add_text(s, x2, y2, w2, Inches(0.6), "UNIT ECONOMICS (Projected)",
             size=14, bold=True, color=WHITE, font="Arial Black",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ue_rows = [
        ("Avg. district size",     "1,000 students"),
        ("Blended price",          "≈ $4 / student / yr"),
        ("Rev. per district",      "$4,000 / yr"),
        ("Gross margin (target)",  "≈ 85%"),
        ("CAC (assumed)",          "$800 / district"),
        ("5-yr LTV (assumed)",     "$18,400"),
        ("LTV : CAC (model)",      "23 : 1"),
    ]
    yy = y2 + Inches(0.85)
    for i, (label, value) in enumerate(ue_rows):
        rr = yy + Inches(i * 0.5)
        add_text(s, x2 + Inches(0.3), rr, Inches(3.0), Inches(0.4),
                 label, size=12, bold=True, color=TEXT)
        add_text(s, x2 + Inches(3.3), rr, w2 - Inches(3.6), Inches(0.4),
                 value, size=14, bold=True, color=FOREST,
                 font="Arial Black", align=PP_ALIGN.RIGHT)

    # RIGHT-TOP: Growth chart
    x3 = Inches(6.85); y3 = Inches(1.7); w3 = Inches(6.1); h3 = Inches(2.7)
    add_rounded(s, x3, y3, w3, h3, WHITE)
    add_text(s, x3 + Inches(0.2), y3 + Inches(0.1), w3 - Inches(0.4), Inches(0.4),
             "3-year district growth (projection)",
             size=14, bold=True, color=FOREST, font="Arial Black")
    add_column_chart(
        s, x3 + Inches(0.2), y3 + Inches(0.55), w3 - Inches(0.4), h3 - Inches(0.7),
        categories=["Y1 pilot", "Y2 regional", "Y3 national"],
        series=[("Districts", [3, 50, 500])],
    )

    # RIGHT-BOTTOM: Business model card
    x4 = Inches(6.85); y4 = Inches(4.55); w4 = Inches(6.1); h4 = Inches(1.75)
    add_rounded(s, x4, y4, w4, h4, CARD_LIGHT)
    add_text(s, x4 + Inches(0.2), y4 + Inches(0.1), w4 - Inches(0.4), Inches(0.4),
             "Business model",
             size=14, bold=True, color=FOREST, font="Arial Black")
    add_bullets(s, x4 + Inches(0.2), y4 + Inches(0.5), w4 - Inches(0.4), h4 - Inches(0.6), [
        "Annual district contracts · Recurring ARR — not ads",
        "Land-and-expand by tier · No data sale, no resale",
    ], size=12)

    add_citation(s, Inches(0.6), Inches(6.9), Inches(12.3),
                 "LTV/CAC is a model output with stated assumptions, not reported results. "
                 "Pricing benchmark on prior slide. 0 paid districts today.")


# ── Slide 10: Demand Generation & Go-To-Market ───────────────────────────────
def slide_gtm(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Demand Generation & Go-To-Market")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Three parallel channels. One decision-maker: the district administrator.",
             size=13, italic=True, color=MUTED)

    channels = [
        ("Design-partner pilots",
         ["Free 6-month pilot",
          "3 design-partner districts",
          "Measured vs. matched control",
          "Case studies → Year 2 collateral"],
         CARD_BLUE),
        ("Conferences + RFPs",
         ["ASU+GSV (Apr · San Diego)",
          "ISTE (Jun · Orlando)",
          "FETC + TCEA (Jan-Feb)",
          "State RFPs: CA + TX first"],
         CARD_SAGE),
        ("Teacher-led pull (PLG)",
         ["Free teacher-tier signup",
          "85% of principals trust teachers (Transcend '24)",
          "MagicSchool reached 6 M teachers this way",
          "→ bottom-up to district contract"],
         DARK_BLOB),
    ]
    x = Inches(0.6)
    y = Inches(1.7)
    w = Inches(4.0)
    h = Inches(3.8)
    gap = Inches(0.15)
    for i, (head, items, color) in enumerate(channels):
        cx = x + (w + gap) * i
        add_rounded(s, cx, y, w, h, WHITE)
        add_rect(s, cx, y, w, Inches(0.6), color)
        add_text(s, cx, y, w, Inches(0.6), head,
                 size=16, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, cx + Inches(0.3), y + Inches(0.8), w - Inches(0.6), h - Inches(0.95),
                    items, size=12)

    # Procurement reality strip
    add_rounded(s, Inches(0.6), Inches(5.7), Inches(12.3), Inches(0.95), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(5.78), Inches(11.9), Inches(0.4),
             "K-12 sales reality — EdWeek Market Brief 2025:",
             size=12, bold=True, color=FOREST, font="Arial Black")
    add_text(s, Inches(0.85), Inches(6.18), Inches(11.9), Inches(0.5),
             "37% of districts close in 6–11 months · Pilot adds ~12 mo · "
             "Title I FY26 $18.5 B + IDEA Part B $15.2 B (post-ESSER, NAESP 2026)",
             size=11, color=TEXT)

    # Audience messaging
    add_rounded(s, Inches(0.6), Inches(6.75), Inches(12.3), Inches(0.65), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.85), Inches(11.9), Inches(0.5),
             "Super: 'Pilot free 6 mo.' · Teacher: 'Grade in mins.' · Parent: 'Weekly digest, plain English.'",
             size=12, bold=True, color=FOREST, italic=True, align=PP_ALIGN.CENTER)


# ── Slide 11: Product Validation / Pilot ─────────────────────────────────────
def slide_pilot(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Product Validation / Pilot")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.45),
             "Pre-launch. 0 paid users. This slide says exactly how we'll know if it works.",
             size=13, italic=True, color=MUTED)

    add_rounded(s, Inches(0.6), Inches(1.75), Inches(5.9), Inches(5.2), WHITE)
    add_rect(s, Inches(0.6), Inches(1.75), Inches(5.9), Inches(0.55), CARD_BLUE)
    add_text(s, Inches(0.6), Inches(1.75), Inches(5.9), Inches(0.55),
             "PILOT DESIGN (FALL 2026)",
             size=14, bold=True, color=WHITE, font="Arial Black",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(4.4), [
        "Target: 3 design-partner districts",
        "Duration: 6 months, free",
        "Controls: matched non-pilot cohort (same district, different grade/section)",
        "Consent: opt-in per-parent, FERPA-aligned",
        "Funding: covered by founding team",
        "Deliverable to partners: per-student progress report + teacher time-use audit",
    ], size=12)

    add_rounded(s, Inches(6.7), Inches(1.75), Inches(6.2), Inches(5.2), WHITE)
    add_rect(s, Inches(6.7), Inches(1.75), Inches(6.2), Inches(0.55), CARD_SAGE)
    add_text(s, Inches(6.7), Inches(1.75), Inches(6.2), Inches(0.55),
             "WHAT WE'LL MEASURE",
             size=14, bold=True, color=WHITE, font="Arial Black",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(6.95), Inches(2.5), Inches(5.7), Inches(4.4), [
        "PRIMARY: NAEP-aligned pre/post assessment delta vs. control",
        "SECONDARY: teacher hrs/week on grading (self-reported log)",
        "SECONDARY: parent weekly-active rate",
        "SECONDARY: student streak days / assignment completion",
        "GATE for Year 2 pricing: positive delta on ≥2 of 3 secondary metrics, ≥2 of 3 districts",
        "Null-result policy: publish negative findings, too",
    ], size=12)

    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.3),
                 "No user-facing outcomes claimed elsewhere in this deck. ESSER expired Sept 2024 — "
                 "pilot funded by founding team, not federal recovery dollars. Null-result publication is policy.")


# ── Slide 15: Competition — Landscape Matrix (colorful 9 vendors) ────────────
def slide_competition_matrix(prs):
    s = new_slide(prs, 0)
    add_page_title(s, "Competition — Landscape")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Nine real products. Vendor pricing pages 2026-04-25.",
             size=13, italic=True, color=MUTED)

    # Header row
    x = Inches(0.6); y = Inches(1.55)
    cols = [Inches(0.22), Inches(2.6), Inches(3.4), Inches(2.5), Inches(1.7), Inches(1.9)]
    headers = ["", "Vendor", "Focus", "$ / student / yr", "Parent", "Modality"]
    cx = x
    for i, hd in enumerate(headers):
        add_rect(s, cx, y, cols[i], Inches(0.45), FOREST)
        if hd:
            add_text(s, cx, y, cols[i], Inches(0.45), hd,
                     size=11, bold=True, color=WHITE, font="Arial Black",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[i]

    # Data rows — Limud first, then by category
    rows = [
        (V_LIMUD,    "Limud",            "Adaptive + 4 roles",          "$4",         "Yes",     "Yes"),
        (V_KHAN,     "Khan / Khanmigo",  "Tutor + content",             "$35",        "Limited", "No"),
        (V_IXL,      "IXL Learning",     "Practice + analytics",        "$5–10",      "Yes",     "No"),
        (V_DREAMBOX, "DreamBox",         "Adaptive math + reading",     "$20–25",     "Limited", "Limited"),
        (V_ALEKS,    "ALEKS",            "Adaptive math",               "Quote only", "No",      "Limited"),
        (V_MAGIC,    "MagicSchool AI",   "Teacher AI assistant",        "~$100 / tch",   "No",      "No"),
        (V_BRISK,    "Brisk Teaching",   "Teacher AI assistant",        "~$100 / tch",   "No",      "No"),
        (V_CURIPOD,  "Curipod",          "AI lesson engagement",        "~$90 / tch",    "Polls",   "No"),
        (V_EDUAIDE,  "Eduaide.AI",       "Teacher AI quizzes",          "~$72 / tch",    "No",      "No"),
    ]
    row_h = Inches(0.42)
    for ri, row in enumerate(rows):
        ry = y + Inches(0.45) + Inches(ri * 0.42)
        is_limud = row[1] == "Limud"
        bg = CARD_SAGE if is_limud else (WHITE if ri % 2 == 0 else CARD_LIGHT)
        cx = x
        # Color stripe (vendor brand color)
        add_rect(s, cx, ry, cols[0], row_h, row[0])
        cx += cols[0]
        for ci, val in enumerate(row[1:]):
            add_rect(s, cx, ry, cols[ci+1], row_h, bg)
            color = WHITE if is_limud else TEXT
            bold = ci == 0 or is_limud
            add_text(s, cx, ry, cols[ci+1], row_h, val,
                     size=11, bold=bold, color=color, font="Calibri",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += cols[ci+1]

    # Two-column segmentation callout
    seg_y = Inches(5.85)
    add_rounded(s, Inches(0.6), seg_y, Inches(6.05), Inches(1.05), CARD_LIGHT)
    add_text(s, Inches(0.85), seg_y + Inches(0.1), Inches(5.6), Inches(0.4),
             "Adaptive student platforms",
             size=13, bold=True, color=CARD_BLUE, font="Arial Black")
    add_text(s, Inches(0.85), seg_y + Inches(0.5), Inches(5.6), Inches(0.5),
             "Khan · IXL · DreamBox · ALEKS — student-facing, district pricing.",
             size=11, color=TEXT)

    add_rounded(s, Inches(6.85), seg_y, Inches(6.05), Inches(1.05), CARD_LIGHT)
    add_text(s, Inches(7.1), seg_y + Inches(0.1), Inches(5.6), Inches(0.4),
             "Teacher AI assistants",
             size=13, bold=True, color=V_MAGIC, font="Arial Black")
    add_text(s, Inches(7.1), seg_y + Inches(0.5), Inches(5.6), Inches(0.5),
             "MagicSchool · Brisk · Curipod · Eduaide — teacher-only, freemium.",
             size=11, color=TEXT)

    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.3),
                 "Sources: ixl.com · khanacademy.org · dreambox.com · aleks.com · "
                 "magicschool.ai · briskteaching.com · curipod.com · eduaide.ai. Accessed 2026-04-25.")


# ── Slide 16: Competition — Funding Heat (NEW) ───────────────────────────────
def slide_competition_funding(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Competition — Funding Heat")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "AI-tutoring VC drew ≈$1.4 B in 2024–25. Recent K-12 / EdTech AI rounds:",
             size=13, italic=True, color=MUTED)

    rounds = [
        ("Preply · Series D",                       150),
        ("Speak · Series C",                         78),
        ("MagicSchool · Series B (Bain)",            45),
        ("Gizmo AI · Series A",                      22),
        ("Brisk Teaching · Series A (Bessemer)",     15),
        ("Curipod · Seed",                            5),
        ("Limud · pre-seed",                          0),
    ]
    add_bar_chart(
        s, Inches(0.6), Inches(1.6), Inches(8.4), Inches(5.0),
        categories=[r[0] for r in rounds],
        values=[r[1] for r in rounds],
    )

    # Right column: callout cards
    x4 = Inches(9.25); y4 = Inches(1.6); w4 = Inches(3.7); h4 = Inches(1.55); gap4 = Inches(0.15)
    callouts = [
        ("$1.4 B",  "AI tutoring VC, 2024–25",     "newmarketpitch sector total",  GOLD),
        ("$45 M",   "MagicSchool Series B",        "Bain Capital · Jan 2025",      V_MAGIC),
        ("$15 M",   "Brisk Teaching Series A",     "Bessemer · Mar 2025",          V_BRISK),
    ]
    for i, (big, lbl, sub, color) in enumerate(callouts):
        cy = y4 + (h4 + gap4) * i
        add_rounded(s, x4, cy, w4, h4, WHITE)
        add_rect(s, x4, cy, Inches(0.18), h4, color)
        add_text(s, x4 + Inches(0.35), cy + Inches(0.1), w4 - Inches(0.5), Inches(0.6),
                 big, size=32, bold=True, color=color, font="Arial Black")
        add_text(s, x4 + Inches(0.35), cy + Inches(0.75), w4 - Inches(0.5), Inches(0.35),
                 lbl, size=12, bold=True, color=FOREST, font="Arial Black")
        add_text(s, x4 + Inches(0.35), cy + Inches(1.1), w4 - Inches(0.5), Inches(0.4),
                 sub, size=10, italic=True, color=MUTED)

    add_citation(s, Inches(0.6), Inches(7.05), Inches(12.3),
                 "Sources: magicschool.ai blog · finsmes.com (Brisk) · techcrunch.com (Gizmo) · "
                 "newmarketpitch.com (sector aggregate). Bars are $ M raised in disclosed rounds.")


# ── Slide 17: Competition — Where Limud Fits (NEW) ───────────────────────────
def slide_competition_differentiators(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Competition — Where Limud Fits")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Three things no other K-12 product does at once.",
             size=13, italic=True, color=MUTED)

    diffs = [
        ("FOUR ROLES",
         "Student · Teacher · Parent · Admin",
         "Khan / IXL / DreamBox / ALEKS = student-only.\n"
         "MagicSchool / Brisk / Curipod / Eduaide = teacher-only.\n"
         "Limud is the only one serving all four in one product.",
         CARD_BLUE),
        ("MODALITY ADAPTS",
         "Visual · Auditory · Kinesthetic",
         "Adaptive comps adjust difficulty + pacing only.\n"
         "Limud re-renders the same lesson per modality —\n"
         "no other K-12 platform claims this in product copy.",
         CARD_SAGE),
        ("TRANSPARENT $",
         "$4 / student blended",
         "MagicSchool / Brisk / ALEKS / DreamBox hide district price.\n"
         "Limud publishes Basic $3 · Pro $5 · Enterprise $8.\n"
         "Budget tier vs. $5–35 adaptive comps.",
         GOLD),
    ]
    x = Inches(0.6); y = Inches(1.65); w = Inches(4.05); h = Inches(4.5); gap = Inches(0.15)
    for i, (head, sub, body, color) in enumerate(diffs):
        cx = x + (w + gap) * i
        add_rounded(s, cx, y, w, h, WHITE)
        add_rect(s, cx, y, w, Inches(1.05), color)
        add_text(s, cx + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
                 head, size=18, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.2), y + Inches(0.65), w - Inches(0.4), Inches(0.4),
                 sub, size=13, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.3), y + Inches(1.3), w - Inches(0.6), h - Inches(1.5),
                 body, size=12, color=TEXT)

    add_rounded(s, Inches(0.6), Inches(6.3), Inches(12.3), Inches(0.85), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.45), Inches(11.9), Inches(0.5),
             "Net: only Limud lets a teacher upload once and a parent get a weekly digest at $4 / student.",
             size=14, bold=True, color=FOREST, font="Arial Black",
             align=PP_ALIGN.CENTER)

    add_citation(s, Inches(0.6), Inches(7.2), Inches(12.3),
                 "Differentiators verified against vendor pricing/feature pages 2026-04-25 (full sources slide 15).")


# ── Slide 18: Competitor Voice (NEW — customer review quotes) ────────────────
def slide_competitor_voice(prs):
    s = new_slide(prs, 0)
    add_page_title(s, "Competitor Voice")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "What teachers and parents actually say. Not marketing copy.",
             size=13, italic=True, color=MUTED)

    quotes = [
        ("KHANMIGO", V_KHAN,
         '"Strong scaffolding for math, but humanities still feels like reading from a textbook."',
         "Common Sense Education teacher · 2025",
         "Strength: Socratic method.\nGap: Subject coverage uneven.",
        ),
        ("IXL", V_IXL,
         '"More a source of stress than education — students sent to guidance counseling over SmartScore."',
         "Trustpilot parent · 2025",
         "Strength: K-12 skill coverage.\nGap: Punitive scoring demotivates.",
        ),
        ("MAGICSCHOOL", V_MAGIC,
         '"Saves me 7-10 hrs/week, but I have to fact-check every output — especially math."',
         "G2 / Common Sense teacher · 2024-25",
         "Strength: Time savings.\nGap: Hallucinations · manual QA tax.",
        ),
    ]
    x = Inches(0.6); y = Inches(1.65); w = Inches(4.05); h = Inches(4.4); gap = Inches(0.15)
    for i, (vendor, color, quote, attribution, takeaway) in enumerate(quotes):
        cx = x + (w + gap) * i
        add_rounded(s, cx, y, w, h, WHITE)
        add_rect(s, cx, y, w, Inches(0.7), color)
        add_text(s, cx, y, w, Inches(0.7), vendor,
                 size=18, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.25), y + Inches(0.95), w - Inches(0.5), Inches(2.0),
                 quote, size=14, italic=True, color=TEXT)
        add_text(s, cx + Inches(0.25), y + Inches(2.95), w - Inches(0.5), Inches(0.4),
                 "— " + attribution, size=10, italic=True, color=MUTED)
        add_rounded(s, cx + Inches(0.2), y + Inches(3.4), w - Inches(0.4), Inches(0.9), CARD_LIGHT)
        add_text(s, cx + Inches(0.35), y + Inches(3.5), w - Inches(0.7), Inches(0.75),
                 takeaway, size=11, color=FOREST)

    add_rounded(s, Inches(0.6), Inches(6.2), Inches(12.3), Inches(0.95), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.3), Inches(11.9), Inches(0.4),
             "Pattern: nobody has trustworthy AI content AND non-punitive engagement.",
             size=14, bold=True, color=FOREST, font="Arial Black")
    add_text(s, Inches(0.85), Inches(6.7), Inches(11.9), Inches(0.4),
             "Limud's product hypothesis: confidence-scored AI + intrinsic motivation, designed for grades 6-12.",
             size=12, italic=True, color=MUTED)

    add_citation(s, Inches(0.6), Inches(7.2), Inches(12.3),
                 "Sources: commonsense.org/education · trustpilot.com/review/ixl.com · g2.com · "
                 "educationnext.org. All accessed 2026-04-27.")


# ── Slide 19: Adjacent Threats (NEW — Big Tech free-tier squeeze) ────────────
def slide_adjacent_threats(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Adjacent Threats")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Free-tier squeeze: Big Tech + incumbent LMS bundling AI for free.",
             size=13, italic=True, color=MUTED)

    R_HIGH = RGBColor(0xC8, 0x42, 0x42)  # red
    R_MED  = RGBColor(0xE0, 0x9A, 0x3F)  # amber
    R_LOW  = RGBColor(0x6B, 0x8E, 0x5E)  # sage

    threats = [
        ("HIGH", R_HIGH, "Google Classroom + Gemini",      "Free for educators · 6M trained · Miami-Dade adopting"),
        ("HIGH", R_HIGH, "Microsoft Copilot Education",    "Bundled M365 · Teams in ~70% of U.S. schools"),
        ("HIGH", R_HIGH, "PowerSchool PowerBuddy",         "50M students · zero-friction in Schoology LMS"),
        ("MED",  R_MED,  "OpenAI ChatGPT for Teachers",    "Free for verified U.S. K-12 teachers thru Jun 2027"),
        ("MED",  R_MED,  "Anthropic Claude for Education", "K-12 via partners (MagicSchool); no standalone yet"),
        ("MED",  R_MED,  "Canvas IgniteAI",                "Free advanced AI thru Jun 2026 (US)"),
        ("LOW",  R_LOW,  "Quizlet Q-Chat / Magic Notes",   "Flashcards + tutor · narrow use case"),
        ("LOW",  R_LOW,  "Duolingo Math",                  "Math-only · K-12 not core"),
    ]
    x = Inches(0.6); y = Inches(1.6)
    cols = [Inches(0.85), Inches(4.6), Inches(6.85)]
    headers = ["Risk", "Threat", "Why it matters"]
    cx = x
    for i, hd in enumerate(headers):
        add_rect(s, cx, y, cols[i], Inches(0.45), FOREST)
        add_text(s, cx, y, cols[i], Inches(0.45), hd,
                 size=11, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[i]
    row_h = Inches(0.5)
    for ri, (risk, rcolor, threat, why) in enumerate(threats):
        ry = y + Inches(0.45) + Inches(ri * 0.5)
        bg = WHITE if ri % 2 == 0 else CARD_LIGHT
        cx = x
        add_rect(s, cx, ry, cols[0], row_h, rcolor)
        add_text(s, cx, ry, cols[0], row_h, risk,
                 size=12, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[0]
        add_rect(s, cx, ry, cols[1], row_h, bg)
        add_text(s, cx + Inches(0.2), ry, cols[1] - Inches(0.2), row_h,
                 threat, size=12, bold=True, color=FOREST, font="Arial Black",
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[1]
        add_rect(s, cx, ry, cols[2], row_h, bg)
        add_text(s, cx + Inches(0.2), ry, cols[2] - Inches(0.2), row_h,
                 why, size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[2]

    # Nightmare scenario callout
    add_rounded(s, Inches(0.6), Inches(6.05), Inches(12.3), Inches(1.1), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.15), Inches(11.9), Inches(0.4),
             "Nightmare scenario — the free-tier squeeze:",
             size=13, bold=True, color=R_HIGH, font="Arial Black")
    add_text(s, Inches(0.85), Inches(6.55), Inches(11.9), Inches(0.55),
             "Google + MS + PowerSchool saturate 'good enough' AI by 2027. Limud's window narrows if "
             "positioning stays generic. Specialize: grades 6-12 modality + parent loop.",
             size=11, color=TEXT)

    add_citation(s, Inches(0.6), Inches(7.2), Inches(12.3),
                 "Sources: blog.google/products-and-platforms/products/education · educatorstechnology.com · "
                 "powerschool.com · openai.com/teachers · claude.com/solutions/education. Accessed 2026-04-27.")


# ── Slide 20: Strategic Plays (NEW — Build / Parity / Monitor / Skip) ────────
def slide_strategic_plays(prs):
    s = new_slide(prs, 2)
    add_page_title(s, "Strategic Plays")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Where to differentiate. Where to match. Where to skip.",
             size=13, italic=True, color=MUTED)

    plays = [
        ("DIFFERENTIATE", "Where we win",
         ["4 roles (no comp does this)",
          "Modality adaptation V/A/K",
          "Confidence-scored AI",
          "Transparent $4 / student"],
         CARD_SAGE),
        ("ACHIEVE PARITY", "What we must match",
         ["Lesson-plan quality vs MagicSchool",
          "Auto-grading speed vs Brisk",
          "Free teacher tier (PLG)",
          "LMS integration: Canvas + Schoology"],
         CARD_BLUE),
        ("MONITOR", "Watch quarterly",
         ["PowerBuddy roadmap (HIGH risk)",
          "Gemini for Education K-12 ships",
          "Copilot Teach pricing post-Dec '25",
          "Anthropic K-12 standalone"],
         GOLD),
        ("DON'T COMPETE", "Skip these lanes",
         ["K-5 elementary (DreamBox)",
          "Math-only practice (ALEKS)",
          "Pure flashcards (Quizlet)",
          "Higher ed (Canvas / OpenAI)"],
         RGBColor(0x6E, 0x7E, 0x8C)),
    ]
    # 2×2 grid
    x0 = Inches(0.6); y0 = Inches(1.65); w = Inches(6.05); h = Inches(2.5); gx = Inches(0.2); gy = Inches(0.18)
    for i, (head, sub, items, color) in enumerate(plays):
        cx = x0 + (w + gx) * (i % 2)
        cy = y0 + (h + gy) * (i // 2)
        add_rounded(s, cx, cy, w, h, WHITE)
        add_rect(s, cx, cy, w, Inches(0.65), color)
        add_text(s, cx + Inches(0.2), cy + Inches(0.05), Inches(3.5), Inches(0.55),
                 head, size=15, bold=True, color=WHITE, font="Arial Black",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(3.5), cy + Inches(0.05), w - Inches(3.7), Inches(0.55),
                 sub, size=11, italic=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, cx + Inches(0.25), cy + Inches(0.8), w - Inches(0.5), h - Inches(0.95),
                    items, size=12)

    add_rounded(s, Inches(0.6), Inches(6.85), Inches(12.3), Inches(0.45), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.92), Inches(11.9), Inches(0.4),
             "Lanes derived from competitor matrix (slide 15) + customer voice (slide 18) + threat map (slide 19).",
             size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 21: Team ───────────────────────────────────────────────────────────
def slide_team(prs):
    s = new_slide(prs, 1)
    add_page_title(s, "Team")
    add_text(s, Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.4),
             "Six builders. Full-stack coverage.",
             size=13, italic=True, color=MUTED)

    members = [
        ("Tamar B",  "CEO",          "Vision, partnerships, district outreach"),
        ("Noam S",   "CTO",          "Architecture, AI pipeline, data model"),
        ("Lior B",   "Design Lead",  "Product design, brand, UX across all four portals"),
        ("Erez O",   "Engineering",  "Student + Parent experiences"),
        ("Eytan B",  "Engineering",  "Teacher tools + AI grading"),
        ("Noam E",   "Engineering",  "Admin + district console"),
    ]
    cols, col_w, col_h = 3, Inches(4.1), Inches(2.1)
    gap_x, gap_y = Inches(0.15), Inches(0.25)
    x0 = Inches(0.5)
    y0 = Inches(1.75)
    for i, (name, role, desc) in enumerate(members):
        cx = x0 + (col_w + gap_x) * (i % cols)
        cy = y0 + (col_h + gap_y) * (i // cols)
        add_rounded(s, cx, cy, col_w, col_h, WHITE)
        add_oval(s, cx + Inches(0.25), cy + Inches(0.3), Inches(1.2), Inches(1.2), CARD_SAGE)
        initial = name[0]
        add_text(s, cx + Inches(0.25), cy + Inches(0.3), Inches(1.2), Inches(1.2),
                 initial, size=36, bold=True, color=WHITE, font="Arial Black",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(1.55), cy + Inches(0.35), col_w - Inches(1.7), Inches(0.4),
                 name, size=16, bold=True, color=FOREST, font="Arial Black")
        add_text(s, cx + Inches(1.55), cy + Inches(0.8), col_w - Inches(1.7), Inches(0.35),
                 role, size=12, bold=True, italic=True, color=CARD_BLUE)
        add_text(s, cx + Inches(1.55), cy + Inches(1.15), col_w - Inches(1.7), Inches(0.85),
                 desc, size=11, color=TEXT)

    add_rounded(s, Inches(0.6), Inches(6.45), Inches(12.3), Inches(0.75), CARD_LIGHT)
    add_text(s, Inches(0.85), Inches(6.55), Inches(11.9), Inches(0.5),
             "We ship in two-week cycles. The rest of this story is ahead of us.",
             size=13, italic=True, color=FOREST)


# ── Slide 14: Thank You / Ask ────────────────────────────────────────────────
def slide_thanks(prs):
    s = new_slide(prs, 2)
    add_logo(s, Inches(6.07), Inches(1.0), size=Inches(1.1))
    add_text(s, Inches(0.6), Inches(2.2), Inches(12.0), Inches(1.0),
             "Thank you.", size=56, bold=True, color=FOREST,
             font="Arial Black", align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.2), Inches(12.0), Inches(0.5),
             "Questions welcome. Skepticism too.",
             size=18, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.66), Inches(3.9), Inches(2.0), Inches(0.04), FOREST)

    # The ask
    add_rounded(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(2.3), WHITE)
    add_rect(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.55), CARD_SAGE)
    add_text(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.55),
             "WHAT WE'RE ASKING FOR",
             size=14, bold=True, color=WHITE, font="Arial Black",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(2.0), Inches(5.05), Inches(9.3), Inches(1.5), [
        "Three design-partner districts for a free 6-month Fall 2026 pilot (grades 6–12 only).",
        "Introductions to middle/high school district superintendents in CA or TX.",
        "Your honest critique of this plan — especially the parts we have wrong.",
    ], size=13)

    add_text(s, Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.4),
             "Every mind learns differently.",
             size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ═════════════════════════════ BUILD ═════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 22-slide structure (2026-04-27 — adds competitive brief slides 18-20):
    slide_title(prs)                          # 1
    slide_framework(prs)                      # 2
    slide_problem_outcomes(prs)               # 3  — Problem split A
    slide_problem_workforce(prs)              # 4  — Problem split B
    slide_solution_approach(prs)              # 5  — Solution split A
    slide_solution_evidence(prs)              # 6  — Solution split B
    slide_market(prs)                         # 7
    slide_workflow_today(prs)                 # 8
    slide_workflow_limud(prs)                 # 9
    slide_use_case(prs)                       # 10
    slide_pricing_tiers(prs)                  # 11 — Pricing split A
    slide_unit_economics_growth(prs)          # 12 — Pricing split B
    slide_gtm(prs)                            # 13
    slide_pilot(prs)                          # 14
    slide_competition_matrix(prs)             # 15 — Competition Landscape
    slide_competition_funding(prs)            # 16 — Competition Funding
    slide_competition_differentiators(prs)    # 17 — Where Limud Fits
    slide_competitor_voice(prs)               # 18 — Customer voice (NEW)
    slide_adjacent_threats(prs)               # 19 — Adjacent threats (NEW)
    slide_strategic_plays(prs)                # 20 — Build/Parity/Monitor (NEW)
    slide_team(prs)                           # 21
    slide_thanks(prs)                         # 22

    out = "Limud_Business_Plan.pptx"
    prs.save(out)
    import os
    print(f"Saved: {os.path.abspath(out)}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
